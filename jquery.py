from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

# Create new PowerPoint presentation
ppt = Presentation()

# Slide data for "jQuery Effects"
slides_data = [
    ("Introduction to jQuery Effects",
     "• jQuery is a lightweight JavaScript library used to simplify HTML document traversal, event handling, and animations.\n"
     "• jQuery effects are pre-defined animations that make web pages interactive and visually appealing.\n"
     "• Common jQuery effects include showing, hiding, fading, sliding, and custom animations.\n"
     "• These effects are created using simple method calls without manually writing complex JavaScript code.\n"
     "• jQuery effects improve user experience by adding smooth transitions and visual feedback."),
    
    ("Basic Show and Hide Effects",
     "• The `show()` and `hide()` methods are used to display or hide HTML elements.\n"
     "• Syntax:\n"
     "   - $(selector).show(speed, callback)\n"
     "   - $(selector).hide(speed, callback)\n"
     "• The `speed` parameter can be 'slow', 'fast', or a time in milliseconds.\n"
     "• The `callback` function runs after the effect completes.\n"
     "• Example:\n"
     "   $(\"#text\").hide(1000);\n"
     "   $(\"#text\").show(1000);\n"
     "• These effects are useful for toggling visibility of elements dynamically."),
    
    ("Fade Effects",
     "• Fading effects gradually change the opacity of elements.\n"
     "• Common methods:\n"
     "   - fadeIn(speed, callback)\n"
     "   - fadeOut(speed, callback)\n"
     "   - fadeToggle(speed, callback)\n"
     "   - fadeTo(speed, opacity, callback)\n"
     "• Example:\n"
     "   $(\"#image\").fadeIn(1500);\n"
     "   $(\"#text\").fadeTo(1000, 0.5);\n"
     "• These effects are ideal for creating smooth appearance and disappearance transitions."),
    
    ("Sliding Effects",
     "• jQuery provides sliding animations that reveal or hide content vertically.\n"
     "• Common methods:\n"
     "   - slideDown(speed, callback)\n"
     "   - slideUp(speed, callback)\n"
     "   - slideToggle(speed, callback)\n"
     "• Example:\n"
     "   $(\"#menu\").slideDown(1000);\n"
     "   $(\"#menu\").slideUp(1000);\n"
     "• These effects are commonly used in dropdown menus, accordions, and collapsible panels."),
    
    ("Custom Animations with animate()",
     "• The `animate()` method allows you to create custom animations by changing CSS properties.\n"
     "• Syntax: $(selector).animate({ properties }, speed, callback);\n"
     "• Example:\n"
     "   $(\"#box\").animate({left: '200px', opacity: '0.5'}, 1000);\n"
     "• You can chain multiple animations together for complex motion effects.\n"
     "• jQuery internally uses CSS transitions to handle these animations smoothly."),
    
    ("Conclusion",
     "• jQuery effects simplify animation handling in web development.\n"
     "• They allow developers to create rich, interactive user interfaces with minimal code.\n"
     "• Common effects include fading, sliding, showing, hiding, and custom animations.\n"
     "• These effects enhance user engagement and provide smooth visual feedback.\n"
     "• With jQuery, developers can implement animations that are both efficient and easy to control.")
]

# Add slides to presentation
for title, content in slides_data:
    slide = ppt.slides.add_slide(ppt.slide_layouts[1])
    slide.shapes.title.text = title
    textbox = slide.placeholders[1]
    textbox.text = content

    # Formatting text
    for para in textbox.text_frame.paragraphs:
        para.font.size = Pt(18)
        para.font.color.rgb = RGBColor(0, 0, 0)

# Save PowerPoint file
output_path = "/mnt/data/jQuery_Effects_Presentation.pptx"
ppt.save(output_path)

output_path
